import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from io import BytesIO
import zipfile
import re
from datetime import datetime

# --- KONFIGURACJA STRONY ---
st.set_page_config(page_title="Generator Katalogu CC", layout="wide", page_icon="📄")

# CSS - Style i Wygląd
st.markdown("""
    <style>
        @import url('https://fonts.googleapis.com/css2?family=Roboto:wght@300;400;700&display=swap');
        html, body, [class*="css"]  {
            font-family: 'Roboto', sans-serif;
        }
        h1 { color: #000000; }
        .stButton>button {
            width: 100%;
            border-radius: 5px;
            height: 3em;
            background-color: #000000;
            color: #ffffff;
            font-weight: bold;
        }
        /* Ciemny motyw obsługiwany jest automatycznie przez Streamlit,
           ale tutaj wymuszamy pewne akcenty dla czytelności */
    </style>
""", unsafe_allow_html=True)

# --- FUNKCJE POMOCNICZE ---

def parse_scale_business(value):
    """
    Zamienia tekst '1 mln PLN', '2,5 mld' na liczbę (float) do sortowania.
    """
    if pd.isna(value):
        return 0.0
    
    text = str(value).lower().replace(',', '.').replace(' ', '')
    multiplier = 1.0
    
    if 'mld' in text or 'b' in text: # miliard
        multiplier = 1_000_000_000.0
    elif 'mln' in text or 'm' in text: # milion
        multiplier = 1_000_000.0
    elif 'tys' in text or 'k' in text: # tysiąc
        multiplier = 1_000.0
        
    # Wyciąganie samej liczby z tekstu
    numbers = re.findall(r"[-+]?\d*\.\d+|\d+", text)
    if numbers:
        return float(numbers[0]) * multiplier
    return 0.0

def clean_polish_typography(text):
    """
    Zapobiega wiszącym spójnikom (wdowy i sieroty).
    Zamienia spację po 'w', 'z', 'i', 'a' na twardą spację.
    """
    if not isinstance(text, str):
        return text
    
    # Lista spójników do przyklejenia
    conjunctions = [" w ", " z ", " i ", " a ", " o ", " u ", " na ", " do "]
    
    for word in conjunctions:
        # Zamiana spacji zwykłej na twardą spację (\u00A0)
        # Regex szuka spójnika otoczonego spacjami, ignorując wielkość liter
        pattern = re.compile(re.escape(word), re.IGNORECASE)
        # Funkcja lambda zachowuje oryginalną wielkość liter spójnika
        text = pattern.sub(lambda m: m.group(0).replace(' ', '\u00A0', 1), text) # zamień tylko drugą spację? 
        # Prościej: po prostu zamieniamy spację PO spójniku na twardą.
        
    # Prosta metoda brute-force dla pewności (zamiast regex dla każdego):
    for word in [" w", " z", " i", " a", " o", " u"]:
        text = text.replace(f"{word} ", f"{word}\u00A0")
        text = text.replace(f"{word.upper()} ", f"{word.upper()}\u00A0")
        
    return text

def fit_text_to_shape(shape, text, max_chars=500):
    """
    Wstawia tekst i zmniejsza czcionkę, jeśli tekst jest długi.
    """
    if not text:
        text = "-"
    
    # Przycięcie tekstu, jeśli za długi (opcja bezpieczeństwa)
    # text = text[:max_chars] if len(text) > max_chars else text
    
    # Wstawienie tekstu
    text_frame = shape.text_frame
    text_frame.clear() # Czyścimy domyślny tekst
    p = text_frame.paragraphs[0]
    p.text = str(text)
    
    # Logika zmniejszania czcionki (prosta heurystyka)
    # Domyślny rozmiar zakładamy, że jest ustawiony w szablonie.
    # Jeśli tekst jest bardzo długi, wymuszamy mniejszy.
    char_count = len(text)
    if char_count > 600:
        p.font.size = Pt(8)
    elif char_count > 400:
        p.font.size = Pt(10)
    elif char_count > 200:
        # Zostawiamy domyślną lub ustawiamy np. 12
        pass 
        
    # Justowanie (opcjonalnie, jeśli szablon tego nie ma)
    # from pptx.enum.text import PP_ALIGN
    # p.alignment = PP_ALIGN.JUSTIFY

def insert_image_into_shape(slide, shape, image_stream):
    """
    Wstawia obrazek w miejsce kształtu (Placeholder lub Shape), zachowując jego wymiary.
    """
    # Pobieramy wymiary i pozycję oryginalnego kształtu (placeholdera)
    left, top = shape.left, shape.top
    width, height = shape.width, shape.height
    
    # Jeśli to Placeholder typu Picture, metoda jest prosta:
    if shape.is_placeholder and shape.placeholder_format.type == 18: # 18 = PICTURE
        try:
            shape.insert_picture(image_stream)
            return
        except:
            pass # Jeśli się nie uda, robimy metodą manualną
            
    # Metoda manualna: dodajemy nowy obrazek na wierzch, stary ukrywamy/usuwamy
    try:
        # Dodajemy zdjęcie
        pic = slide.shapes.add_picture(image_stream, left, top, width, height)
        
        # Próbujemy usunąć stary kształt (placeholder) - w python-pptx to nie zawsze działa idealnie
        # więc często po prostu przykrywamy go nowym zdjęciem.
        # shape.element.delete() 
    except Exception as e:
        print(f"Błąd wstawiania obrazka: {e}")

def create_gray_placeholder(slide, shape):
    """
    Rysuje szary prostokąt w miejscu brakującego zdjęcia.
    """
    left, top = shape.left, shape.top
    width, height = shape.width, shape.height
    
    # Dodajemy kształt
    from pptx.enum.shapes import MSO_SHAPE
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(200, 200, 200) # Szary
    rect.line.color.rgb = RGBColor(150, 150, 150)
    
    # Opcjonalnie tekst "BRAK ZDJĘCIA"
    tf = rect.text_frame
    tf.text = "BRAK PLIKU"
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

# --- GŁÓWNA APLIKACJA ---

st.title("Generator Katalogu CC 🚀")
st.markdown("Aplikacja przekształca dane z Excela w gotową prezentację PowerPoint.")

# 1. SIDEBAR - USTAWIENIA I PLIKI
with st.sidebar:
    st.header("1. Wgraj pliki")
    uploaded_excel = st.file_uploader("Baza Danych (Excel/CSV)", type=['xlsx', 'csv'])
    uploaded_pptx = st.file_uploader("Szablon Slajdu (.pptx)", type=['pptx'])
    uploaded_zip = st.file_uploader("Zdjęcia i Logotypy (.zip)", type=['zip'], help="Folder ze zdjęciami spakowany do ZIP.")
    
    st.info("Pamiętaj: Nazwy plików w kolumnie 'Photo' i 'Logo' w Excelu muszą pasować do plików w ZIP.")

# 2. LOGIKA PRZETWARZANIA
if uploaded_excel and uploaded_pptx:
    # Wczytywanie danych
    try:
        if uploaded_excel.name.endswith('.csv'):
            df = pd.read_csv(uploaded_excel)
        else:
            df = pd.read_excel(uploaded_excel)
            
        # Standaryzacja nazw kolumn (usuwamy białe znaki)
        df.columns = df.columns.str.strip()
        
        # Wymagane kolumny (mapowanie nazw z Twojego pliku)
        # Dostosuj te nazwy jeśli w Excelu są inne!
        col_map = {
            "Imię": "Imię",
            "Nazwisko": "Nazwisko",
            "Firma": "Firma",
            "Branża": "Branża",
            "Skala Biznesu": "Skala Biznesu", # lub "Skala biznesu"
            "Grupa CC": "Grupa CC", # lub "Grupa"
            "Photo": "Photo nazwa pliku", # Nazwa pliku zdjęcia w Excelu
            "Logo": "Logo nazwa pliku",   # Nazwa pliku logo w Excelu
            "Opis": "Katalog Członków CC - opis do 500 znaków"
        }
        
        # Sprawdzamy czy kolumny istnieją (case insensitive search)
        df_cols_lower = [c.lower() for c in df.columns]
        
        def get_real_col_name(target):
            for col in df.columns:
                if col.lower() == target.lower():
                    return col
            # Fallback dla specyficznych nazw z Twoich plików CSV
            if target == "Photo nazwa pliku" and "Photo" in df.columns: return "Photo" # Czasami po prostu Photo
            return target

        # Aktualizujemy mapowanie o prawdziwe nazwy z pliku
        real_col_map = {k: get_real_col_name(v) for k, v in col_map.items()}
        
        # Obliczenie wartości numerycznej dla sortowania
        sort_col_name = real_col_map["Skala Biznesu"]
        if sort_col_name in df.columns:
            df["_sort_value"] = df[sort_col_name].apply(parse_scale_business)
        else:
            df["_sort_value"] = 0

        # --- SEKCJA GŁÓWNA ---
        
        st.divider()
        st.subheader("2. Wybierz i Posortuj Rekordy")
        
        col1, col2 = st.columns(2)
        with col1:
            # Filtrowanie Grup
            group_col = real_col_map["Grupa CC"]
            if group_col in df.columns:
                all_groups = df[group_col].dropna().unique().tolist()
                selected_groups = st.multiselect("Wybierz Grupy CC:", all_groups, default=all_groups)
                df_filtered = df[df[group_col].isin(selected_groups)].copy()
            else:
                df_filtered = df.copy()
                
        with col2:
            # Sortowanie
            sort_option = st.selectbox(
                "Sortowanie:",
                ["Domyślne (jak w pliku)", "Skala Biznesu (Największa -> Najmniejsza)", "Nazwisko (A-Z)", "Firma (A-Z)"]
            )
            
            if sort_option == "Skala Biznesu (Największa -> Najmniejsza)":
                df_filtered = df_filtered.sort_values(by="_sort_value", ascending=False)
            elif sort_option == "Nazwisko (A-Z)":
                df_filtered = df_filtered.sort_values(by=real_col_map["Nazwisko"], ascending=True)
            elif sort_option == "Firma (A-Z)":
                df_filtered = df_filtered.sort_values(by=real_col_map["Firma"], ascending=True)

        # Tabela do wyboru
        df_filtered.insert(0, "Wybierz", True)
        
        edited_df = st.data_editor(
            df_filtered,
            column_config={
                "Wybierz": st.column_config.CheckboxColumn("Generuj?", default=True),
                "_sort_value": None # Ukrywamy kolumnę techniczną
            },
            hide_index=True,
            use_container_width=True,
            height=400
        )
        
        final_selection = edited_df[edited_df["Wybierz"] == True]
        st.markdown(f"**Liczba wybranych slajdów:** {len(final_selection)}")

        # Obsługa ZIP (case insensitive mapping)
        images_map = {} # 'nazwapliku.jpg' -> bytes
        if uploaded_zip:
            with zipfile.ZipFile(uploaded_zip) as z:
                for filename in z.namelist():
                    # Ignorujemy foldery (__MACOSX to śmieci z Maca)
                    if not filename.endswith('/') and '__MACOSX' not in filename:
                        # Kluczem jest nazwa pliku małymi literami (bez ścieżki folderów)
                        simple_name = filename.split('/')[-1].lower()
                        images_map[simple_name] = z.read(filename)

        # --- GENEROWANIE ---
        st.divider()
        if st.button("GENERUJ PREZENTACJĘ ⚡", type="primary"):
            if len(final_selection) == 0:
                st.error("Wybierz przynajmniej jeden rekord!")
            else:
                progress_bar = st.progress(0)
                status_text = st.empty()
                
                # Ładowanie szablonu
                prs = Presentation(uploaded_pptx)
                # Zakładamy, że pierwszy slajd to wzorzec. Będziemy go powielać.
                # W python-pptx najlepiej używać slide_layout. 
                # Jeśli szablon ma 1 slajd, jego layout to zwykle layouty przypisane do Mastera.
                # Użyjemy layoutu pierwszego slajdu jako bazy.
                base_layout = prs.slides[0].slide_layout 
                
                # Usuwamy slajd wzorcowy z finalnej prezentacji (opcjonalnie, zależy czy user chce)
                # Na razie zostawiamy, żeby user widział, potem usuwamy?
                # Najlepiej: tworzymy nowe slajdy, a na końcu usuwamy pierwszy (wzorcowy).
                
                total_slides = len(final_selection)
                
                for i, (index, row) in enumerate(final_selection.iterrows()):
                    # Aktualizacja paska postępu
                    progress = (i + 1) / total_slides
                    progress_bar.progress(progress)
                    status_text.text(f"Generowanie slajdu {i+1}/{total_slides}: {row.get(real_col_map['Nazwisko'], '')}")
                    
                    # 1. Dodaj nowy slajd
                    slide = prs.slides.add_slide(base_layout)
                    
                    # 2. Iterujemy po wszystkich kształtach na nowym slajdzie
                    # Szukamy placeholderów i tekstów do podmiany
                    for shape in slide.placeholders:
                        shape_name = shape.name.upper() # np. "TXT_IMIE" lub domyślne nazwy PPT
                        shape_text = shape.text if hasattr(shape, 'text') else ""
                        
                        # --- MAPOWANIE TEKSTÓW ---
                        # Sprawdzamy czy nazwa kształtu LUB tekst w nim zawiera klucz
                        
                        # Funkcja do bezpiecznego pobierania wartości
                        def get_val(key):
                            val = row.get(real_col_map[key], "-")
                            return val if pd.notna(val) else "-"

                        # Imię
                        if "{IMIĘ}" in shape_text or "IMIĘ" in shape_name or "IMIE" in shape_name:
                            fit_text_to_shape(shape, get_val("Imię"))
                        
                        # Nazwisko
                        elif "{NAZWISKO}" in shape_text or "NAZWISKO" in shape_name:
                            fit_text_to_shape(shape, get_val("Nazwisko"))
                            
                        # Firma
                        elif "{FIRMA}" in shape_text or "FIRMA" in shape_name:
                            fit_text_to_shape(shape, get_val("Firma"))
                            
                        # Branża
                        elif "{BRANŻA}" in shape_text or "BRANŻA" in shape_name or "BRANZA" in shape_name:
                            fit_text_to_shape(shape, get_val("Branża"))
                            
                        # Skala
                        elif "{SKALA" in shape_text or "SKALA" in shape_name:
                            fit_text_to_shape(shape, get_val("Skala Biznesu"))
                        
                        # Grupa
                        elif "{GRUPA" in shape_text or "GRUPA" in shape_name:
                            fit_text_to_shape(shape, get_val("Grupa CC"))
                            
                        # Opis (z poprawą typografii)
                        elif "{OPIS}" in shape_text or "OPIS" in shape_name or "500" in shape_text:
                            raw_desc = get_val("Opis")
                            clean_desc = clean_polish_typography(str(raw_desc))
                            fit_text_to_shape(shape, clean_desc)

                        # --- MAPOWANIE ZDJĘĆ ---
                        
                        # Zdjęcie (PHOTO)
                        elif "PHOTO" in shape_name or "ZDJĘCIE" in shape_name:
                            photo_file = str(row.get(real_col_map["Photo"], "")).lower().strip()
                            if uploaded_zip and photo_file in images_map:
                                insert_image_into_shape(slide, shape, BytesIO(images_map[photo_file]))
                            else:
                                create_gray_placeholder(slide, shape)
                        
                        # Logo (LOGO)
                        elif "LOGO" in shape_name:
                            logo_file = str(row.get(real_col_map["Logo"], "")).lower().strip()
                            if uploaded_zip and logo_file in images_map:
                                insert_image_into_shape(slide, shape, BytesIO(images_map[logo_file]))
                            else:
                                # create_gray_placeholder(slide, shape) # Dla logo można zostawić puste
                                pass

                    # Obsługa statycznych tekstów (nie-placeholderów), jeśli chcemy podmieniać np. "{Imię}" w zwykłym polu tekstowym
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    # Prosta podmiana w tekście ciągłym
                                    if "{Imię}" in run.text:
                                        run.text = run.text.replace("{Imię}", str(row.get(real_col_map["Imię"], "")))
                                    if "{Nazwisko}" in run.text:
                                        run.text = run.text.replace("{Nazwisko}", str(row.get(real_col_map["Nazwisko"], "")))
                                    # itd...

                # Usuwamy pierwszy slajd (wzorzec) - hack dla python-pptx
                # (dostęp do wewnętrznej listy slajdów XML)
                xml_slides = prs.slides._sldIdLst
                slides = list(xml_slides)
                xml_slides.remove(slides[0])

                # Zapis
                output = BytesIO()
                prs.save(output)
                output.seek(0)
                
                # Nazwa pliku
                timestamp = datetime.now().strftime("%Y.%m.%d_%H:%M")
                file_name = f"Katalog_CC_Generated_{timestamp}.pptx"
                
                status_text.success("Gotowe! ✅")
                st.download_button(
                    label="POBIERZ PREZENTACJĘ",
                    data=output,
                    file_name=file_name,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    type="primary"
                )

    except Exception as e:
        st.error(f"Wystąpił błąd podczas przetwarzania: {e}")
        st.write("Szczegóły błędu:", e)
else:
    # Ekran powitalny
    st.info("👈 Wgraj pliki w pasku bocznym, aby rozpocząć.")
